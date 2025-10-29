# docintel.py — Robust Azure Document Intelligence client (analyze + polling + fallback)
import time
import json
import requests
import streamlit as st
from typing import Optional, List
from config import CONFIG
import io, os, mimetypes

_HEADERS_BIN = {
    "Ocp-Apim-Subscription-Key": CONFIG["AI_DOC_INTEL_KEY"],
    # binary 업로드시 Content-Type은 호출 시 지정
}
_HEADERS_JSON = {
    "Ocp-Apim-Subscription-Key": CONFIG["AI_DOC_INTEL_KEY"],
    "Content-Type": "application/json"
}

def _endpoint() -> str:
    ep = CONFIG["AI_DOC_INTEL_ENDPOINT"].rstrip("/")
    if not ep.startswith("https://"):
        ep = "https://" + ep
    return ep

def _poll_operation_result(op_location: str, api_version: str, timeout_sec: int = 60, interval: float = 1.2) -> dict:
    deadline = time.time() + timeout_sec
    headers = {"Ocp-Apim-Subscription-Key": CONFIG["AI_DOC_INTEL_KEY"]}
    while time.time() < deadline:
        r = requests.get(op_location, headers=headers, timeout=30)
        r.raise_for_status()
        data = r.json()
        status = data.get("status") or data.get("status", "").lower()
        if status in ("succeeded", "Succeeded"):
            return data
        if status in ("failed", "Failed"):
            raise RuntimeError(f"DocIntel analyze failed: {json.dumps(data, ensure_ascii=False)[:800]}")
        time.sleep(interval)
    raise TimeoutError("DocIntel analyze polling timeout")

def _analyze_bytes(content: bytes, mime_type: str, model: str, api_version: str) -> dict:
    url = f"{_endpoint()}/documentintelligence/documentModels/{model}:analyze?_overload=analyzeDocument&api-version={api_version}"
    headers = dict(_HEADERS_BIN)
    headers["Content-Type"] = mime_type or "application/octet-stream"

    # 시작(비동기) → 202 + operation-location
    r = requests.post(url, headers=headers, data=content, timeout=60)
    if r.status_code == 404:
        # 모델/버전/경로 불일치일 확률 큼 — 상위에서 폴백 기회 제공
        raise requests.HTTPError("404 Not Found", response=r)
    r.raise_for_status()

    op_loc = r.headers.get("operation-location") or r.headers.get("Operation-Location")
    if not op_loc:
        # 최신 버전은 반드시 LRO를 반환
        raise RuntimeError(f"operation-location header missing: {r.text[:500]}")

    # 폴링
    result = _poll_operation_result(op_loc, api_version=api_version)
    return result

def _extract_text_from_result(result: dict) -> str:
    # 2024-07-31 / 2023-07-31 공통 구조: analyzeResult → content / pages/lines/words 등
    ar = result.get("analyzeResult") or {}
    # 1) 전체 content 있으면 가장 간단
    content = ar.get("content")
    if content:
        return content

    # 2) fallback: 페이지-라인을 이어붙임
    pages: List[dict] = ar.get("pages") or []
    out_lines = []
    for p in pages:
        lines = p.get("lines") or []
        for ln in lines:
            txt = ln.get("content")
            if txt:
                out_lines.append(txt)
    return "\n".join(out_lines).strip()

def extract_text_docintel(content: bytes, mime_type: Optional[str] = None) -> str:
    """
    Azure Document Intelligence로 텍스트 추출 (비동기 폴링 + 다중 모델/버전 폴백)
    - 기본: 2024-11-30 + prebuilt-read
    - 폴백: 2024-11-30 + prebuilt-layout → prebuilt-read → prebuilt-layout
    """
    preferred_model = (CONFIG.get("AI_DOC_INTEL_MODEL") or "prebuilt-read").strip() or "prebuilt-read"
    api_versions = ["2024-11-30"]
    models = [preferred_model, "prebuilt-read", "prebuilt-layout"]

    last_err = None
    for ver in api_versions:
        for model in models:
            try:
                res = _analyze_bytes(content, mime_type or "application/octet-stream", model=model, api_version=ver)
                txt = _extract_text_from_result(res)
                if not txt:
                    raise RuntimeError("Empty text extracted")
                return txt
            except requests.HTTPError as he:
                # 404면 다음 조합으로 시도
                if he.response is not None and he.response.status_code == 404:
                    # 상세 메시지 화면 표시
                    try:
                        detail = he.response.json()
                    except Exception:
                        detail = he.response.text
                    st.info(f"DocIntel 404 (ver={ver}, model={model}) → 다음 조합 시도\n{detail}")
                    last_err = he
                    continue
                # 그 외 HTTP 에러는 즉시 표시 후 다음 조합
                st.warning(f"DocIntel HTTPError (ver={ver}, model={model}): {he}")
                last_err = he
                continue
            except Exception as e:
                # 네트워크/타임아웃 등
                st.warning(f"DocIntel Exception (ver={ver}, model={model}): {e}")
                last_err = e
                continue

    # 모든 조합 실패
    if last_err:
        raise last_err
    raise RuntimeError("DocIntel analyze failed with no additional details")


# 선택 의존성: 없는 경우에도 graceful fallback
try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import chardet
except Exception:
    chardet = None


def _guess_ext(filename: str) -> str:
    return os.path.splitext(filename or "")[-1].lower()


def _detect_text_encoding(data: bytes) -> str:
    if not data:
        return "utf-8"
    if chardet:
        try:
            enc = chardet.detect(data).get("encoding") or "utf-8"
            return enc
        except Exception:
            pass
    # 안전 기본값
    return "utf-8"


def _extract_txt_like(filename: str, data: bytes) -> str:
    enc = _detect_text_encoding(data)
    try:
        return data.decode(enc, errors="replace")
    except Exception:
        # 최후의 보루
        return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    if not PyPDF2:
        return "[PDF parser(PyPDF2) 미설치] requirements.txt에 PyPDF2 추가하세요."
    try:
        buf = io.BytesIO(data)
        reader = PyPDF2.PdfReader(buf)
        texts = []
        for page in reader.pages:
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t:
                texts.append(t)
        return "\n\n".join(texts).strip() or "[빈 PDF이거나 텍스트 추출 실패]"
    except Exception as e:
        return f"[PDF 추출 오류] {e}"


def _extract_docx(data: bytes) -> str:
    if not docx:
        return "[DOCX parser(python-docx) 미설치] requirements.txt에 python-docx 추가하세요."
    try:
        buf = io.BytesIO(data)
        d = docx.Document(buf)
        paras = [p.text for p in d.paragraphs if p.text]
        return "\n".join(paras).strip() or "[DOCX 본문이 비어있음]"
    except Exception as e:
        return f"[DOCX 추출 오류] {e}"


def _extract_pptx(data: bytes) -> str:
    if not Presentation:
        return "[PPTX parser(python-pptx) 미설치] requirements.txt에 python-pptx 추가하세요."
    try:
        buf = io.BytesIO(data)
        prs = Presentation(buf)
        lines = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines).strip() or "[PPTX 텍스트가 없습니다]"
    except Exception as e:
        return f"[PPTX 추출 오류] {e}"


def _extract_xlsx(data: bytes) -> str:
    if not openpyxl:
        return "[XLSX parser(openpyxl) 미설치] requirements.txt에 openpyxl 추가하세요."
    try:
        buf = io.BytesIO(data)
        wb = openpyxl.load_workbook(buf, data_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"[Sheet] {ws.title}")
            for row in ws.iter_rows(values_only=True):
                row_vals = ["" if v is None else str(v) for v in row]
                lines.append("\t".join(row_vals))
        return "\n".join(lines).strip() or "[XLSX 내용이 없습니다]"
    except Exception as e:
        return f"[XLSX 추출 오류] {e}"


def extract_text_naive(filename: str, content: bytes) -> str:
    """
    간단 텍스트 추출기 (로컬 파싱)
    - PDF: PyPDF2
    - DOCX: python-docx
    - PPTX: python-pptx
    - XLSX: openpyxl (탭으로 셀 연결)
    - TXT/MD/기타 텍스트: chardet로 인코딩 추정 후 디코드
    - 그 외: MIME/확장자에 따라 best-effort
    """
    ext = _guess_ext(filename)

    # 확장자 우선 분기
    if ext == ".pdf":
        return _extract_pdf(content)
    if ext == ".docx":
        return _extract_docx(content)
    if ext == ".pptx":
        return _extract_pptx(content)
    if ext in (".xlsx", ".xlsm"):
        return _extract_xlsx(content)
    if ext in (".txt", ".md", ".csv", ".log"):
        return _extract_txt_like(filename, content)

    # 확장자가 애매할 때 MIME 힌트 시도
    mt, _ = mimetypes.guess_type(filename or "")
    if mt:
        if mt == "application/pdf":
            return _extract_pdf(content)
        if mt in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",):
            return _extract_docx(content)
        if mt in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",):
            return _extract_pptx(content)
        if mt in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",):
            return _extract_xlsx(content)
        if mt.startswith("text/"):
            return _extract_txt_like(filename, content)

    # 마지막 fallback: 텍스트 시도
    return _extract_txt_like(filename, content)
